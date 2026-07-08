import fitz
import pandas as pd
import pytesseract

from PIL import Image
from docx import Document
from pptx import Presentation


# Optional OCR path (Windows)
# Uncomment if needed
# pytesseract.pytesseract.tesseract_cmd = (
#     r"C:\Program Files\Tesseract-OCR\tesseract.exe"
# )


def extract_file_content(file_name, content):
    """
    Extract text from uploaded files.

    Supported:
    - PDF
    - DOCX
    - TXT
    - CSV
    - XLSX
    - PPTX
    - JPG
    - JPEG
    - PNG
    """

    try:

        content.seek(0)

        ext = file_name.lower().split(".")[-1]

        # =========================
        # PDF
        # =========================
        if ext == "pdf":

            text = ""

            pdf = fitz.open(
                stream=content.read(),
                filetype="pdf"
            )

            for page in pdf:
                text += page.get_text()

            return text.strip()

        # =========================
        # DOCX
        # =========================
        elif ext == "docx":

            content.seek(0)

            doc = Document(content)

            text = "\n".join(
                paragraph.text
                for paragraph in doc.paragraphs
            )

            return text.strip()

        # =========================
        # TXT
        # =========================
        elif ext == "txt":

            content.seek(0)

            return (
                content.read()
                .decode("utf-8", errors="ignore")
                .strip()
            )

        # =========================
        # CSV
        # =========================
        elif ext == "csv":

            content.seek(0)

            return (
                content.read()
                .decode("utf-8", errors="ignore")
                .strip()
            )

        # =========================
        # XLSX
        # =========================
        elif ext == "xlsx":

            content.seek(0)

            df = pd.read_excel(content)

            return df.to_string()

        # =========================
        # PPTX
        # =========================
        elif ext == "pptx":

            content.seek(0)

            prs = Presentation(content)

            text = ""

            for slide in prs.slides:

                for shape in slide.shapes:

                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

            return text.strip()

        # =========================
        # IMAGES
        # =========================
        elif ext in [
            "jpg",
            "jpeg",
            "png"
        ]:

            content.seek(0)

            image = Image.open(content)

            text = pytesseract.image_to_string(
                image
            )

            return text.strip()

        # =========================
        # UNSUPPORTED
        # =========================
        return f"Unsupported file type: {ext}"

    except Exception as e:

        return (
            f"FILE_PARSE_ERROR: "
            f"{type(e).__name__}: {str(e)}"
        )